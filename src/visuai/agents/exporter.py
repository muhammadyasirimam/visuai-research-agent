"""
Exporter Agent - Handles export to multiple formats.

Author: Muhammad Yasir Imam
"""

import os
from datetime import datetime

from ..core.config import Config


class ExporterAgent:
    """Agent responsible for exporting results to various formats."""

    def __init__(self, config: Config):
        self.config = config

    def export(self, charts, narrative, formats, session_id):
        """Export results to specified formats."""
        export_paths = {}
        output_dir = f"./visuai_output/{session_id}"
        os.makedirs(output_dir, exist_ok=True)
        for fmt in formats:
            try:
                if fmt == "html":
                    path = self._export_html(charts, narrative, output_dir, session_id)
                    export_paths["html"] = path
                elif fmt == "pdf":
                    path = self._export_pdf(charts, narrative, output_dir, session_id)
                    export_paths["pdf"] = path
                elif fmt == "png":
                    path = self._export_png(charts, output_dir, session_id)
                    export_paths["png"] = path
                elif fmt == "jupyter":
                    path = self._export_jupyter(charts, narrative, output_dir, session_id)
                    export_paths["jupyter"] = path
                elif fmt == "latex":
                    path = self._export_latex(charts, narrative, output_dir, session_id)
                    export_paths["latex"] = path
                elif fmt == "pptx":
                    path = self._export_pptx(charts, narrative, output_dir, session_id)
                    export_paths["pptx"] = path
                elif fmt == "markdown":
                    path = self._export_markdown(charts, narrative, output_dir, session_id)
                    export_paths["markdown"] = path
            except Exception as e:
                print(f"   ⚠️ Export to {fmt} failed: {e}")
        return export_paths

    def _export_html(self, charts, narrative, output_dir, session_id):
        html_content = self._generate_html_report(charts, narrative, session_id)
        path = os.path.join(output_dir, f"report_{session_id}.html")
        with open(path, "w", encoding="utf-8") as f:
            f.write(html_content)
        return path

    def _generate_html_report(self, charts, narrative, session_id):
        chart_sections = []
        for chart in charts:
            img_tag = ""
            if chart.figure_path and os.path.exists(chart.figure_path):
                import base64
                with open(chart.figure_path, "rb") as img_file:
                    img_data = base64.b64encode(img_file.read()).decode()
                img_tag = f'<img src="data:image/png;base64,{img_data}" alt="{chart.title}" style="max-width:100%;border-radius:8px;box-shadow:0 4px 12px rgba(0,0,0,0.1);"/>'
            insights_html = ""
            if chart.insights:
                insights_html = "<ul>" + "".join([f"<li>{ins}</li>" for ins in chart.insights]) + "</ul>"
            section = f'<div class="chart-card"><h3>{chart.title}</h3><p><strong>Type:</strong> {chart.chart_type}</p>{img_tag}<div class="insights"><h4>Insights</h4>{insights_html}</div></div>'
            chart_sections.append(section)
        charts_html = "\n".join(chart_sections)
        dt = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>VisuAI Report - {session_id}</title>
<style>
:root {{ --primary: #2563eb; --secondary: #7c3aed; --accent: #06b6d4; --bg: #0f172a; --card: #1e293b; --text: #f1f5f9; --muted: #94a3b8; }}
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ font-family: "Segoe UI", system-ui, sans-serif; background: var(--bg); color: var(--text); line-height: 1.6; }}
.container {{ max-width: 1400px; margin: 0 auto; padding: 2rem; }}
.header {{ text-align: center; padding: 3rem 0; background: linear-gradient(135deg, var(--primary), var(--secondary)); border-radius: 16px; margin-bottom: 2rem; }}
.header h1 {{ font-size: 2.5rem; margin-bottom: 0.5rem; }}
.header p {{ color: rgba(255,255,255,0.8); }}
.narrative {{ background: var(--card); padding: 2rem; border-radius: 12px; margin-bottom: 2rem; border-left: 4px solid var(--accent); }}
.chart-card {{ background: var(--card); padding: 2rem; border-radius: 12px; margin-bottom: 2rem; box-shadow: 0 4px 20px rgba(0,0,0,0.2); }}
.chart-card h3 {{ color: var(--accent); margin-bottom: 1rem; }}
.chart-card img {{ margin: 1rem 0; }}
.insights {{ background: rgba(6,182,212,0.1); padding: 1rem; border-radius: 8px; margin-top: 1rem; }}
.insights h4 {{ color: var(--accent); margin-bottom: 0.5rem; }}
.insights ul {{ padding-left: 1.5rem; }}
.insights li {{ margin-bottom: 0.3rem; }}
.footer {{ text-align: center; padding: 2rem; color: var(--muted); margin-top: 3rem; border-top: 1px solid rgba(255,255,255,0.1); }}
@media print {{ body {{ background: white; color: black; }} .header {{ background: #2563eb; }} }}
</style>
</head>
<body>
<div class="container">
<div class="header">
<h1>🔬 VisuAI Research Report</h1>

    def _export_pdf(self, charts, narrative, output_dir, session_id):
        try:
            from weasyprint import HTML
            html_path = self._export_html(charts, narrative, output_dir, session_id)
            pdf_path = os.path.join(output_dir, f"report_{session_id}.pdf")
            HTML(html_path).write_pdf(pdf_path)
            return pdf_path
        except Exception as e:
            print(f"   ⚠️ PDF export failed: {e}")
            return None

    def _export_png(self, charts, output_dir, session_id):
        png_dir = os.path.join(output_dir, "png")
        os.makedirs(png_dir, exist_ok=True)
        for chart in charts:
            if chart.figure_path and os.path.exists(chart.figure_path):
                import shutil
                shutil.copy2(chart.figure_path, os.path.join(png_dir, f"{chart.chart_id}.png"))
        return png_dir

    def _export_jupyter(self, charts, narrative, output_dir, session_id):
        try:
            import nbformat as nbf
            nb = nbf.v4.new_notebook()
            nb.cells.append(nbf.v4.new_markdown_cell("# VisuAI Research Report"))
            nb.cells.append(nbf.v4.new_markdown_cell("## Executive Summary"))
            for chart in charts:
                if chart.figure_path and os.path.exists(chart.figure_path):
                    code = f"from IPython.display import Image, display"
                    nb.cells.append(nbf.v4.new_code_cell(code))
            path = os.path.join(output_dir, f"report_{session_id}.ipynb")
            with open(path, "w") as f2:
                nbf.write(nb, f2)
            return path
        except Exception as e:
            print(f"   ⚠️ Jupyter export failed: {e}")
            return None

    def _export_latex(self, charts, narrative, output_dir, session_id):
        latex = "\\documentclass[11pt,a4paper]{article}"
        latex += "\\usepackage[utf8]{inputenc}"
        latex += "\\usepackage{graphicx}"
        latex += "\\title{VisuAI Research Report}"
        latex += "\\author{Muhammad Yasir Imam}"
        latex += "\\begin{document}"
        latex += "\\maketitle"
        latex += "\\section{Executive Summary}"
        latex += narrative
        latex += "\\section{Visualizations}"
        for chart in charts:
            if chart.figure_path and os.path.exists(chart.figure_path):
                latex += f"\\subsection{{{chart.title}}}"
        latex += "\\end{document}"
        path = os.path.join(output_dir, f"report_{session_id}.tex")
        with open(path, "w", encoding="utf-8") as f2:
            f2.write(latex)
        return path

    def _export_pptx(self, charts, narrative, output_dir, session_id):
        try:
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "VisuAI Research Report"
            for chart in charts:
                if chart.figure_path and os.path.exists(chart.figure_path):
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = chart.title
                    slide.shapes.add_picture(chart.figure_path, Inches(1), Inches(1.5), width=Inches(8))
            path = os.path.join(output_dir, f"report_{session_id}.pptx")
            prs.save(path)
            return path
        except Exception as e:
            print(f"   ⚠️ PPTX export failed: {e}")
            return None

    def _export_markdown(self, charts, narrative, output_dir, session_id):
        md = "# VisuAI Research Report\n\n"
        md += f"**Session:** {session_id}\n\n"
        md += "## Executive Summary\n\n"
        md += narrative + "\n\n"
        md += "## Visualizations\n\n"
        for chart in charts:
            md += f"### {chart.title}\n\n"
            md += f"- **Type:** {chart.chart_type}\n\n"
        md += "---\n\n*Generated by VisuAI Research Agent v1.0.0*"
        path = os.path.join(output_dir, f"report_{session_id}.md")
        with open(path, "w", encoding="utf-8") as f2:
            f2.write(md)
        return path
